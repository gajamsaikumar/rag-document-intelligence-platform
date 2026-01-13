
from __future__ import annotations#to remove the circular imports
import os
import sys
import json
import hashlib
import shutil
import uuid
from pathlib import Path
from typing import Iterable,List,Optional,Dict,Any

import fitz  # PyMuPDF
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.storage import LocalFileStore
from langchain.retrievers.multi_vector import MultiVectorRetriever

from utils.model_loader import ModelLoader
from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException

from utils.file_io import _session_id, save_uploaded_files
from utils.document_ops import load_documents
from src.document_chat.multimodal.indexer import SessionPaths,create_session_paths,partition_pdf_to_modalities,summarize_texts,summarize_tables,summarize_images
from logger import GLOBAL_LOGGER as log

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx", ".md", ".csv", ".xlsx", ".xls", ".db", ".sqlite", ".sqlite3"}

class FaissManager:

    def __init__(self,index_dir:Path,model_loader:Optional[ModelLoader]=None):
        self.index_dir = Path(index_dir)
        self.index_dir.mkdir(parents=True, exist_ok=True)
        
        self.meta_path = self.index_dir / "ingested_meta.json"
        self._meta: Dict[str, Any] = {"rows": {}}
        
        if self.meta_path.exists():
            try:
                self._meta = json.loads(self.meta_path.read_text(encoding="utf-8")) or {"rows": {}}
            except Exception:
                self._meta = {"rows": {}}
        

        self.model_loader = model_loader or ModelLoader()
        self.emb = self.model_loader.load_embeddings()
        self.vs: Optional[FAISS] = None

    def _exists(self):
        return (self.index_dir / "index.faiss").exists() and (self.index_dir / "index.pkl").exists()
    
    @staticmethod
    def _fingerprint(text: str, md: Dict[str, Any]):
        """
        This function is used to doesnot contains any duplicate info from the faiss database.
        """
        src = md.get("source") or md.get("file_path")
        rid = md.get("row_id")
        if src is not None:
            return f"{src}::{'' if rid is None else rid}"
        return hashlib.sha256(text.encode("utf-8")).hexdigest()

    def _save_meta(self):
        self.meta_path.write_text(json.dumps(self._meta, ensure_ascii=False, indent=2), encoding="utf-8")


    def load_or_create(self,texts:Optional[List[str]]=None, metadatas: Optional[List[dict]] = None):
        if self._exists():
            self.vs = FAISS.load_local(
                str(self.index_dir),
                embeddings=self.emb,
                allow_dangerous_deserialization=True,
            )
            return self.vs
        if not texts:
            raise DocumentPortalException("No existing FAISS index and no data to create one", sys)
        
        self.vs = FAISS.from_texts(texts=texts, embedding=self.emb, metadatas=metadatas or [])
        self.vs.save_local(str(self.index_dir))
        return self.vs

    def add_documents(self,docs: List[Document]):
        if self.vs is None:
            raise RuntimeError("Call load_or_create() before add_documents_idempotent().")
        
        new_docs: List[Document] = []
        
        for d in docs:
            key = self._fingerprint(d.page_content, d.metadata or {})
            if key in self._meta["rows"]:
                continue
            self._meta["rows"][key] = True
            new_docs.append(d)
            
        if new_docs:
            self.vs.add_documents(new_docs)
            self.vs.save_local(str(self.index_dir))
            self._save_meta()
        return len(new_docs)


class DocHandler:
    def __init__(self,data_dir: Optional[str] = None, session_id: Optional[str] = None):
        self.data_dir = data_dir or os.getenv("DATA_STORAGE_PATH", os.path.join(os.getcwd(), "data", "document_analysis"))
        self.session_id = session_id or _session_id("session")
        self.session_path = os.path.join(self.data_dir, self.session_id)
        os.makedirs(self.session_path, exist_ok=True)
        log.info("DocHandler initialized", session_id=self.session_id, session_path=self.session_path)
    
    # -------- Generic multi-format save --------
    def save_file(self, uploaded_file) -> str:
        try:
            filename = os.path.basename(uploaded_file.name)
            ext = os.path.splitext(filename)[1].lower()
            allowed = {".pdf", ".docx", ".pptx", ".md", ".txt", ".xlsx", ".xls", ".csv", ".db", ".sqlite", ".sqlite3"}
            if ext not in allowed:
                raise ValueError(f"Unsupported file type: {ext}. Allowed: {sorted(allowed)}")
            save_path = os.path.join(self.session_path, filename)
            with open(save_path, "wb") as f:
                if hasattr(uploaded_file, "read"):
                    f.write(uploaded_file.read())
                else:
                    f.write(uploaded_file.getbuffer())
            log.info("File saved successfully", file=filename, save_path=save_path, session_id=self.session_id)
            return save_path
        except Exception as e:
            log.error("Failed to save file", error=str(e), session_id=self.session_id)
            raise DocumentPortalException(f"Failed to save file: {str(e)}", e) from e

    # -------- Generic multi-format read --------
    def read_text(self, path: str) -> str:
        try:
            ext = os.path.splitext(path)[1].lower()
            if ext == ".pdf":
                return self.read_pdf(path)
            if ext == ".docx":
                return self._read_docx(path)
            if ext == ".pptx":
                return self._read_pptx(path)
            if ext == ".md":
                return self._read_md(path)
            if ext == ".txt":
                return self._read_txt(path)
            if ext == ".csv":
                return self._read_csv(path)
            if ext == ".xlsx":
                return self._read_xlsx(path)
            if ext == ".xls":
                return self._read_xls(path)
            if ext in {".db", ".sqlite", ".sqlite3"}:
                return self._read_sqlite(path)
            raise ValueError(f"Unsupported extension for reading: {ext}")
        except Exception as e:
            log.error("Failed to read file", error=str(e), file_path=path, session_id=self.session_id)
            raise DocumentPortalException(f"Could not process file: {path}", e) from e
    
    def _read_docx(self, path: str) -> str:
        import docx2txt
        try:
            text = docx2txt.process(path) or ""
            log.info("DOCX read successfully", file_path=path)
            return text
        except Exception as e:
            log.error("Failed to read DOCX", error=str(e), file_path=path)
            raise

    def _read_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    parts.append(f"\n--- Slide {slide_idx} ---\n" + "\n".join(slide_text))
            text = "\n".join(parts)
            log.info("PPTX read successfully", file_path=path, slides=len(prs.slides))
            return text
        except Exception as e:
            log.error("Failed to read PPTX", error=str(e), file_path=path)
            raise

    def _read_md(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            log.info("MD read successfully", file_path=path)
            return content
        except Exception as e:
            log.error("Failed to read MD", error=str(e), file_path=path)
            raise

    def _read_txt(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            log.info("TXT read successfully", file_path=path)
            return content
        except Exception as e:
            log.error("Failed to read TXT", error=str(e), file_path=path)
            raise

    def _read_csv(self, path: str) -> str:
        import csv
        try:
            lines = []
            with open(path, "r", encoding="utf-8", newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    lines.append(", ".join("" if c is None else str(c) for c in row))
            text = "\n".join(lines)
            log.info("CSV read successfully", file_path=path, rows=len(lines))
            return text
        except Exception as e:
            log.error("Failed to read CSV", error=str(e), file_path=path)
            raise

    def _read_xlsx(self, path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"\n--- Sheet: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            text = "\n".join(parts)
            log.info("XLSX read successfully", file_path=path, sheets=len(wb.worksheets))
            return text
        except Exception as e:
            log.error("Failed to read XLSX", error=str(e), file_path=path)
            raise

    def _read_xls(self, path: str) -> str:
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            parts = []
            for sheet in wb.sheets():
                parts.append(f"\n--- Sheet: {sheet.name} ---")
                for rx in range(sheet.nrows):
                    row = [sheet.cell_value(rx, cx) for cx in range(sheet.ncols)]
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            text = "\n".join(parts)
            log.info("XLS read successfully", file_path=path, sheets=wb.nsheets)
            return text
        except Exception as e:
            log.error("Failed to read XLS", error=str(e), file_path=path)
            raise

    def _read_sqlite(self, path: str) -> str:
        try:
            import sqlite3
            # Open database read-only
            uri = f"file:{path}?mode=ro"
            conn = sqlite3.connect(uri, uri=True)
            conn.row_factory = sqlite3.Row
            cur = conn.cursor()

            # List tables
            cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%' ORDER BY name;")
            tables = [r[0] for r in cur.fetchall()]
            if not tables:
                log.info("SQLite DB has no user tables", file_path=path)
                return ""

            parts: list[str] = []
            for t in tables:
                parts.append(f"\n--- Table: {t} ---")
                # Get columns
                try:
                    cur.execute(f"PRAGMA table_info('{t}')")
                    cols = [row[1] for row in cur.fetchall()]
                except Exception:
                    cols = []
                if cols:
                    parts.append("# Columns: " + ", ".join(cols))
                # Dump limited rows
                try:
                    cur.execute(f"SELECT * FROM '{t}' LIMIT 1000")
                    rows = cur.fetchall()
                    for r in rows:
                        if isinstance(r, sqlite3.Row):
                            vals = [r[k] for k in r.keys()]
                        else:
                            vals = list(r)
                        parts.append("\t".join("" if v is None else str(v) for v in vals))
                except Exception as e:  # pragma: no cover - best-effort
                    parts.append(f"# Error reading table '{t}': {e}")

            conn.close()
            text = "\n".join(parts)
            log.info("SQLite read successfully", file_path=path, tables=len(tables))
            return text
        except Exception as e:
            log.error("Failed to read SQLite DB", error=str(e), file_path=path)
            raise


    def save_pdf(self,uploaded_file):
        try:
            filename = os.path.basename(uploaded_file.name)
            if not filename.lower().endswith(".pdf"):
                raise ValueError("Invalid file type. Only PDFs are allowed.")
            save_path = os.path.join(self.session_path, filename)
            with open(save_path, "wb") as f:
                if hasattr(uploaded_file, "read"):
                    f.write(uploaded_file.read())
                else:
                    f.write(uploaded_file.getbuffer())
            log.info("PDF saved successfully", file=filename, save_path=save_path, session_id=self.session_id)
            return save_path
        except Exception as e:
            log.error("Failed to save PDF", error=str(e), session_id=self.session_id)
            raise DocumentPortalException(f"Failed to save PDF: {str(e)}", e) from e

    def read_pdf(self,pdf_path: str):
        try:
            text_chunks = []
            with fitz.open(pdf_path) as doc:
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    text_chunks.append(f"\n--- Page {page_num + 1} ---\n{page.get_text()}")  # type: ignore
            text = "\n".join(text_chunks)
            log.info("PDF read successfully", pdf_path=pdf_path, session_id=self.session_id, pages=len(text_chunks))
            return text
        except Exception as e:
            log.error("Failed to read PDF", error=str(e), pdf_path=pdf_path, session_id=self.session_id)
            raise DocumentPortalException(f"Could not process PDF: {pdf_path}", e) from e

class DocCompare:
    ALLOWED_EXTS = {".pdf", ".docx", ".pptx", ".md", ".txt", ".xlsx", ".xls", ".csv", ".db", ".sqlite", ".sqlite3"}
    def __init__(self,base_dir: str = "data/document_compare", session_id: Optional[str] = None):
        self.base_dir = Path(base_dir)
        self.session_id = session_id or _session_id()
        self.session_path = self.base_dir / self.session_id
        self.session_path.mkdir(parents=True, exist_ok=True)
        log.info("DocumentComparator initialized", session_path=str(self.session_path))
    
    def save_uploaded_files(self,reference_file,actual_file):
        try:
            ref_path = self.session_path / reference_file.name
            act_path = self.session_path / actual_file.name
            for fobj, out in ((reference_file, ref_path), (actual_file, act_path)):
                if not fobj.name.lower().endswith(".pdf"):
                    raise ValueError("Only PDF files are allowed.")
                with open(out, "wb") as f:
                    if hasattr(fobj, "read"):
                        f.write(fobj.read())
                    else:
                        f.write(fobj.getbuffer())
            log.info("Files saved", reference=str(ref_path), actual=str(act_path), session=self.session_id)
            return ref_path, act_path
        except Exception as e:
            log.error("Error saving PDF files", error=str(e), session=self.session_id)
            raise DocumentPortalException("Error saving files", e) from e

    def read_pdf(self,pdf_path: Path):
        try:
            with fitz.open(pdf_path) as doc:
                if doc.is_encrypted:
                    raise ValueError(f"PDF is encrypted: {pdf_path.name}")
                parts = []
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    text = page.get_text()  # type: ignore
                    if text.strip():
                        parts.append(f"\n --- Page {page_num + 1} --- \n{text}")
            log.info("PDF read successfully", file=str(pdf_path), pages=len(parts))
            return "\n".join(parts)
        except Exception as e:
            log.error("Error reading PDF", file=str(pdf_path), error=str(e))
            raise DocumentPortalException("Error reading PDF", e) from e
    
    def read_document(self, file_path: Path) -> str:
        """
        Generic multi-format reader that reuses DocHandler implementations for non-PDF files.
        """
        try:
            ext = file_path.suffix.lower()
            if ext == ".pdf":
                return self.read_pdf(file_path)
            dh = DocHandler(data_dir=str(self.base_dir), session_id=self.session_id)
            return dh.read_text(str(file_path))
        except Exception as e:
            log.error("Error reading document", file=str(file_path), error=str(e))
            raise DocumentPortalException("Error reading document", e) from e

    def combine_docs(self):
        try:
            doc_parts = []
            for file in sorted(self.session_path.iterdir()):
                if file.is_file() and file.suffix.lower() in self.ALLOWED_EXTS:
                    content = self.read_document(file)
                    if content.strip():
                        doc_parts.append(f"Document: {file.name}\n{content}")
            combined_text = "\n\n".join(doc_parts)
            log.info("Documents combined", count=len(doc_parts), session=self.session_id)
            return combined_text
        except Exception as e:
            log.error("Error combining documents", error=str(e), session=self.session_id)
            raise DocumentPortalException("Error combining documents", e) from e

    def clean_old_session(self,keep_latest: int = 3):
        try:
            sessions = sorted([f for f in self.base_dir.iterdir() if f.is_dir()], reverse=True)
            for folder in sessions[keep_latest:]:
                shutil.rmtree(folder, ignore_errors=True)
                log.info("Old session folder deleted", path=str(folder))
        except Exception as e:
            log.error("Error cleaning old sessions", error=str(e))
            raise DocumentPortalException("Error cleaning old sessions", e) from e

class ChatIngestor:
    def __init__(self,temp_base: str = "data",
        faiss_base: str = "faiss_index",
        use_session_dirs: bool = True,
        session_id: Optional[str] = None,):
        try:
            log = CustomLogger().get_logger(__name__)
            self.model_loader = ModelLoader()
            
            self.use_session = use_session_dirs
            self.session_id = session_id or _session_id()
            
            self.temp_base = Path(temp_base); self.temp_base.mkdir(parents=True, exist_ok=True)
            self.faiss_base = Path(faiss_base); self.faiss_base.mkdir(parents=True, exist_ok=True)
            
            self.temp_dir = self._resolve_dir(self.temp_base)
            self.faiss_dir = self._resolve_dir(self.faiss_base)
            
            log.info("ChatIngestor initialized",
                          session_id=self.session_id,
                          temp_dir=str(self.temp_dir),
                          faiss_dir=str(self.faiss_dir),
                          sessionized=self.use_session)
        except Exception as e:
            log.error("Failed to initialize ChatIngestor", error=str(e))
            raise DocumentPortalException("Initialization error in ChatIngestor", e) from e
    
    def _resolve_dir(self,base: Path):
        if self.use_session:
            d = base / self.session_id
            d.mkdir(parents=True, exist_ok=True)
            return d
        return base
    
    def _split(self,docs: List[Document], chunk_size=1000, chunk_overlap=200):
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        chunks = splitter.split_documents(docs)
        log.info("Documents split", chunks=len(chunks), chunk_size=chunk_size, overlap=chunk_overlap)
        return chunks

    def built_retriever(self,uploaded_files: Iterable,
        *,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        k: int = 5,):
        try:
            paths = save_uploaded_files(uploaded_files, self.temp_dir)
            docs = load_documents(paths)
            if not docs:
                raise ValueError("No valid documents loaded")
            
            chunks = self._split(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
            fm = FaissManager(self.faiss_dir, self.model_loader)
            
            texts = [c.page_content for c in chunks]
            metas = [c.metadata for c in chunks]
            
            try:
                vs = fm.load_or_create(texts=texts, metadatas=metas)
            except Exception:
                vs = fm.load_or_create(texts=texts, metadatas=metas)
                
            added = fm.add_documents(chunks)
            log.info("FAISS index updated", added=added, index=str(self.faiss_dir))
            
            return vs.as_retriever(search_type="similarity", search_kwargs={"k": k})
            
        except Exception as e:
            log.error("Failed to build retriever", error=str(e))
            raise DocumentPortalException("Failed to build retriever", e) from e


class MultiModalChatIngestor:
    def __init__(self,temp_base: str = "data", faiss_base: str = "faiss_index", use_session_dirs: bool = True, session_id: str | None = None):

        try:
            self.model_loader = ModelLoader()
            self.use_session = use_session_dirs
            self.session_id = session_id or _session_id()

            self.temp_base = Path(temp_base); self.temp_base.mkdir(parents=True, exist_ok=True)
            self.faiss_base = Path(faiss_base); self.faiss_base.mkdir(parents=True, exist_ok=True)

            sp = create_session_paths(self.faiss_base, self.session_id, self.temp_base)
            self.session_paths: SessionPaths = sp

            cfg = self.model_loader.config.get("multimodal", {})
            self.hires_strategy = bool(cfg.get("hires_strategy", True))
            self.max_images = int(cfg.get("max_images", 8))

            log.info("MultiModalChatIngestor initialized",
                     session_id=self.session_id,
                     faiss_dir=str(self.session_paths.faiss_dir),
                     mm_store=str(self.session_paths.mm_store_dir),
                     assets=str(self.session_paths.assets_dir))
        except Exception as e:
            log.error("Failed to initialize MultiModalChatIngestor", error=str(e))
            raise DocumentPortalException("Initialization error in MultiModalChatIngestor", sys)

    def _load_or_create_vectorstore(self, docs: List[Document]) -> FAISS:
        emb = self.model_loader.load_embeddings()
        index_dir = self.session_paths.faiss_dir
        index_path_faiss = index_dir / "index.faiss"
        index_path_pkl = index_dir / "index.pkl"

        if index_path_faiss.exists() and index_path_pkl.exists():
            vs = FAISS.load_local(str(index_dir), embeddings=emb, allow_dangerous_deserialization=True)
            if docs:
                vs.add_documents(docs)
                vs.save_local(str(index_dir))
            return vs

        # Create fresh from documents
        if not docs:
            raise DocumentPortalException("No documents to create FAISS index", sys)
        vs = FAISS.from_documents(docs, emb)
        vs.save_local(str(index_dir))
        return vs

    def built_retriver(self, uploaded_files: Iterable, *, k: int = 5):
        try:
            # Save uploads to session temp dir
            paths = save_uploaded_files(uploaded_files, self.session_paths.assets_dir.parent)  # keep under session data dir

            texts_all: List[str] = []
            tables_all: List[str] = []
            images_all: List[Path] = []

            for p in paths:
                if p.suffix.lower() == ".pdf":
                    out = partition_pdf_to_modalities(
                        p,
                        self.session_paths.assets_dir,
                        hires_strategy=self.hires_strategy,
                        max_images=self.max_images,
                    )
                    texts_all.extend(out.get("texts", []))
                    tables_all.extend(out.get("tables", []))
                    images_all.extend(out.get("images", []))
                else:
                    log.warning("Skipping non-PDF for multimodal ingestion", path=str(p))

            # Summarize per modality
            text_summaries = summarize_texts(texts_all, self.model_loader)
            table_summaries = summarize_tables(tables_all, self.model_loader)
            image_summaries = summarize_images(images_all, self.model_loader)

            # Build summary docs and payload mappings
            id_key = "doc_id"
            docs: List[Document] = []
            kv_pairs: List[tuple[str, str]] = []

            for i, s in enumerate(text_summaries):
                doc_id = str(uuid.uuid4())
                docs.append(Document(page_content=s, metadata={id_key: doc_id, "modality": "text"}))
                kv_pairs.append((doc_id, texts_all[i]))

            for i, s in enumerate(table_summaries):
                doc_id = str(uuid.uuid4())
                docs.append(Document(page_content=s, metadata={id_key: doc_id, "modality": "table"}))
                kv_pairs.append((doc_id, tables_all[i]))

            for i, s in enumerate(image_summaries):
                doc_id = str(uuid.uuid4())
                docs.append(Document(page_content=s, metadata={id_key: doc_id, "modality": "image"}))
                kv_pairs.append((doc_id, str(images_all[i])))

            # Vectorstore (FAISS) and store (LocalFileStore)
            vs = self._load_or_create_vectorstore(docs)
            store = LocalFileStore(str(self.session_paths.mm_store_dir))
            if kv_pairs:
                # LocalFileStore expects bytes values; encode UTF-8 for strings/paths
                kv_bytes = []
                for k, v in kv_pairs:
                    if isinstance(v, bytes):
                        kv_bytes.append((k, v))
                    else:
                        kv_bytes.append((k, str(v).encode("utf-8")))
                store.mset(kv_bytes)

            retriever = MultiVectorRetriever(vectorstore=vs, docstore=store, id_key=id_key)
            # Return MultiVectorRetriever; k will be applied at query time via search_kwargs
            retriever.search_kwargs = {"k": k}
            return retriever

        except Exception as e:
            log.error("Failed to build multimodal retriever", error=str(e))
            raise DocumentPortalException("Failed to build multimodal retriever", e) from e


